import json
import subprocess
import sys
import tempfile
from pathlib import Path


def read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    blocks = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if text:
                blocks.append(text)
    return "\n".join(blocks)


def read_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    text = "\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()
    if text:
        return text
    with tempfile.TemporaryDirectory(prefix="medvoice-pdf-ocr-") as directory:
        prefix = Path(directory) / "page"
        subprocess.run(
            ["pdftoppm", "-f", "1", "-l", "12", "-png", "-r", "160", str(path), str(prefix)],
            check=True,
            capture_output=True,
        )
        pages = sorted(Path(directory).glob("page-*.png"))
        return "\n\n".join(read_image(page) for page in pages).strip()


def read_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                if text:
                    texts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"[Slide {index}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def run_text_command(command: list[str]) -> str:
    completed = subprocess.run(command, check=True, capture_output=True)
    for encoding in ("utf-8", "gb18030", "latin-1"):
        try:
            return completed.stdout.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    return completed.stdout.decode("utf-8", errors="ignore").strip()


def read_image(path: Path) -> str:
    # Tesseract writes UTF-8 to stdout. Chinese and English are enabled in the
    # production image; fall back to English-only when the Chinese pack is absent.
    try:
        return run_text_command(["tesseract", str(path), "stdout", "-l", "chi_sim+eng"])
    except (subprocess.CalledProcessError, FileNotFoundError):
        return run_text_command(["tesseract", str(path), "stdout", "-l", "eng"])


def main() -> None:
    path = Path(sys.argv[1])
    suffix = path.suffix.lower()
    if suffix == ".docx":
        text = read_docx(path)
    elif suffix == ".doc":
        text = run_text_command(["antiword", str(path)])
    elif suffix == ".pdf":
        text = read_pdf(path)
    elif suffix == ".pptx":
        text = read_pptx(path)
    elif suffix == ".ppt":
        text = run_text_command(["catppt", str(path)])
    elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}:
        text = read_image(path)
    else:
        text = path.read_text(encoding="utf-8", errors="ignore")
    print(json.dumps({"text": text}, ensure_ascii=False))


if __name__ == "__main__":
    main()
